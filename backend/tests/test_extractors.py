import pytest
from app.extractors.pdf_extractor import extract_pdf
from app.extractors.docx_extractor import extract_docx
from app.extractors.excel_extractor import extract_excel
from app.extractors.pptx_extractor import extract_pptx
from app.extractors.ocr_extractor import extract_image_ocr
from app.extractors.html_extractor import extract_html
from app.extractors.router import extract_document_content


def test_pdf_extractor_happy_path(tmp_path):
    """Test PDF extraction on a sample PDF document."""
    pdf_path = tmp_path / "sample.pdf"
    pdf_bytes = b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj 2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj 3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Resources<<>>>>endobj\nxref\n0 4\n0000000000 65535 f\n0000000009 00000 n\n0000000052 00000 n\n0000000101 00000 n\ntrailer<</Size 4/Root 1 0 R>>\nstartxref\n178\n%%EOF"
    pdf_path.write_bytes(pdf_bytes)

    res = extract_pdf(str(pdf_path))
    assert "pages" in res
    assert res["page_count"] >= 1


def test_pdf_extractor_corrupted_file(tmp_path):
    """Test PDF extraction failure path on corrupted file raises ValueError."""
    corrupt_pdf = tmp_path / "corrupt.pdf"
    corrupt_pdf.write_bytes(b"NOT_A_REAL_PDF_DATA_GARBAGE")
    with pytest.raises(Exception):
        extract_pdf(str(corrupt_pdf))


def test_docx_extractor_happy_path(tmp_path):
    """Test DOCX extraction happy path."""
    import docx
    docx_path = tmp_path / "sample.docx"
    doc = docx.Document()
    doc.add_heading("SmartDoc Contract", level=1)
    doc.add_paragraph("This is a test paragraph for DOCX extraction verification.")
    doc.save(str(docx_path))

    res = extract_docx(str(docx_path))
    assert res["page_count"] >= 1
    assert len(res["pages"]) >= 1
    assert "SmartDoc Contract" in res["pages"][0]["text"]


def test_docx_extractor_corrupted_file(tmp_path):
    """Test DOCX extraction failure path on corrupted file raises Exception."""
    corrupt_docx = tmp_path / "corrupt.docx"
    corrupt_docx.write_bytes(b"CORRUPTED_DOCX_BYTES")
    with pytest.raises(Exception):
        extract_docx(str(corrupt_docx))


def test_excel_extractor_happy_path(tmp_path):
    """Test XLSX tabular extraction happy path."""
    import pandas as pd
    xlsx_path = tmp_path / "sample.xlsx"
    df = pd.DataFrame({
        "Revenue": [1000, 2000, 3000],
        "Quarter": ["Q1", "Q2", "Q3"]
    })
    df.to_excel(str(xlsx_path), index=False)

    res = extract_excel(str(xlsx_path), "sample.xlsx")
    assert res["is_tabular"] is True
    assert "tabular_stats" in res
    assert len(res["pages"]) >= 1


def test_excel_extractor_corrupted_file(tmp_path):
    """Test XLSX extraction failure path on invalid bytes raises Exception."""
    corrupt_xlsx = tmp_path / "corrupt.xlsx"
    corrupt_xlsx.write_bytes(b"INVALID_EXCEL_BYTES")
    with pytest.raises(Exception):
        extract_excel(str(corrupt_xlsx), "corrupt.xlsx")


def test_pptx_extractor_happy_path(tmp_path):
    """Test PPTX extraction happy path."""
    from pptx import Presentation
    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Quarterly Business Review"
    prs.save(str(pptx_path))

    res = extract_pptx(str(pptx_path))
    assert res["page_count"] >= 1
    assert "Quarterly Business Review" in res["pages"][0]["text"]


def test_pptx_extractor_corrupted_file(tmp_path):
    """Test PPTX extraction failure path on corrupted file raises Exception."""
    corrupt_pptx = tmp_path / "corrupt.pptx"
    corrupt_pptx.write_bytes(b"CORRUPTED_PPTX_BYTES")
    with pytest.raises(Exception):
        extract_pptx(str(corrupt_pptx))


def test_html_extractor_happy_path(tmp_path):
    """Test HTML extraction happy path."""
    html_path = tmp_path / "sample.html"
    html_path.write_text("<html><body><h1>SmartDoc Overview</h1><p>Web scraper content test.</p></body></html>")

    res = extract_html(str(html_path))
    assert res["page_count"] >= 1
    assert "SmartDoc Overview" in res["pages"][0]["text"]


def test_html_extractor_corrupted_file(tmp_path):
    """Test HTML extraction handling on empty or malformed file."""
    corrupt_html = tmp_path / "corrupt.html"
    corrupt_html.write_text("")
    res = extract_html(str(corrupt_html))
    assert "pages" in res


def test_ocr_extractor_happy_path(tmp_path):
    """Test OCR extractor with small PIL synthetic image."""
    from PIL import Image, ImageDraw
    img_path = tmp_path / "test_ocr.png"
    img = Image.new("RGB", (200, 50), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((10, 10), "TEST", fill=(0, 0, 0))
    img.save(str(img_path))

    res = extract_image_ocr(str(img_path))
    assert "pages" in res
    assert res["page_count"] >= 1


def test_extractor_router(tmp_path):
    """Test extraction router selects correct extractor based on extension."""
    txt_path = tmp_path / "readme.txt"
    txt_path.write_text("Plain text content for router verification.")
    res = extract_document_content(str(txt_path), "readme.txt")
    assert res["file_type"] == "text"
    assert "Plain text content" in res["pages"][0]["text"]
