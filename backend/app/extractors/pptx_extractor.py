import pptx
from typing import Dict, Any, List

def extract_pptx(file_path: str) -> Dict[str, Any]:
    pages = []
    full_text_parts = []
    
    try:
        prs = pptx.Presentation(file_path)
        for idx, slide in enumerate(prs.slides):
            slide_text_parts = []
            
            # Slide shapes text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            slide_text_parts.append(txt)
                            
            # Slide notes text if available
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text_parts.append(f"[Speaker Notes]: {notes}")
                    
            slide_content = "\n".join(slide_text_parts) if slide_text_parts else "[Empty Slide]"
            pages.append({"page_number": idx + 1, "text": slide_content})
            full_text_parts.append(f"--- Slide {idx + 1} ---\n{slide_content}")
            
    except Exception as e:
        raise ValueError(f"Failed to process PowerPoint presentation: {str(e)}")

    full_text = "\n\n".join(full_text_parts) if full_text_parts else "Empty Presentation"
    
    return {
        "file_type": "pptx",
        "page_count": len(pages),
        "full_text": full_text,
        "pages": pages,
        "is_tabular": False,
        "tabular_stats": {}
    }
